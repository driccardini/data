from __future__ import annotations

from collections import Counter
from dataclasses import dataclass
from pathlib import Path
import re
from typing import List

import streamlit as st
from pptx import Presentation


STOPWORDS = {
    "de",
    "la",
    "el",
    "y",
    "en",
    "los",
    "las",
    "del",
    "con",
    "para",
    "una",
    "un",
    "por",
    "que",
    "se",
    "al",
    "a",
    "es",
    "como",
    "más",
    "mas",
    "su",
    "sus",
    "o",
}

KEYWORDS_RIESGO = {
    "riesgo",
    "caída",
    "caida",
    "retraso",
    "coste",
    "costo",
    "deuda",
    "incidencia",
    "brecha",
    "pérdida",
    "perdida",
}

KEYWORDS_OPORTUNIDAD = {
    "crecimiento",
    "ahorro",
    "optimización",
    "optimizacion",
    "eficiencia",
    "expansión",
    "expansion",
    "margen",
    "ingreso",
    "ingresos",
    "automatización",
    "automatizacion",
}


@dataclass
class SlideInfo:
    number: int
    title: str
    body: str


def extraer_slides(ruta_pptx: Path) -> List[SlideInfo]:
    prs = Presentation(str(ruta_pptx))
    slides: List[SlideInfo] = []

    for idx, slide in enumerate(prs.slides, start=1):
        textos = []
        title = ""

        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = (shape.text or "").strip()
            if text:
                textos.append(text)

        body = "\n".join(dict.fromkeys(textos))
        slides.append(SlideInfo(number=idx, title=title or f"Slide {idx}", body=body))

    return slides


def tokenizar(texto: str) -> List[str]:
    palabras = re.findall(r"[A-Za-zÁÉÍÓÚáéíóúÑñ0-9]{3,}", texto.lower())
    return [w for w in palabras if w not in STOPWORDS]


def top_temas(slides: List[SlideInfo], top_n: int = 8) -> List[str]:
    corpus = "\n".join(s.body for s in slides)
    tokens = tokenizar(corpus)
    return [palabra for palabra, _ in Counter(tokens).most_common(top_n)]


def detectar_items(slides: List[SlideInfo], keywords: set[str], max_items: int = 4) -> List[str]:
    hallazgos = []
    for slide in slides:
        texto = f"{slide.title} {slide.body}".lower()
        if any(k in texto for k in keywords):
            resumen = slide.body.split("\n")[0][:140].strip()
            hallazgos.append(f"Slide {slide.number}: {slide.title} — {resumen}")
    return hallazgos[:max_items]


def construir_guion(slides: List[SlideInfo], tono: str, minutos: int) -> str:
    temas = top_temas(slides)
    riesgos = detectar_items(slides, KEYWORDS_RIESGO)
    oportunidades = detectar_items(slides, KEYWORDS_OPORTUNIDAD)

    intro = {
        "Directo": "Voy al punto con los hallazgos críticos y decisiones propuestas.",
        "Narrativo": "Te cuento la historia de 2025: qué pasó, qué aprendimos y qué haremos ahora.",
        "Data-first": "Empiezo por métricas y tendencias para fundamentar cada decisión.",
    }[tono]

    bloques = [
        "# Guion ejecutivo (asistido por IA)",
        f"Duración objetivo: {minutos} minutos",
        "",
        "## Apertura (30-45s)",
        intro,
        "",
        "## Mensajes clave",
    ]

    for i, tema in enumerate(temas[:5], start=1):
        bloques.append(f"{i}. Enfatizar `{tema}` como eje estratégico.")

    bloques.extend(["", "## Riesgos a vigilar"])
    if riesgos:
        bloques.extend([f"- {r}" for r in riesgos])
    else:
        bloques.append("- No se detectaron riesgos explícitos por palabra clave; revisar cualitativamente.")

    bloques.extend(["", "## Oportunidades de impacto"])
    if oportunidades:
        bloques.extend([f"- {o}" for o in oportunidades])
    else:
        bloques.append("- No se detectaron oportunidades explícitas por palabra clave; ampliar con contexto de negocio.")

    bloques.extend(
        [
            "",
            "## Cierre (30s)",
            "Solicitar 3 acuerdos: prioridad del próximo trimestre, dueño por iniciativa, y fecha de checkpoint.",
        ]
    )

    return "\n".join(bloques)


def mensaje_clave(temas: List[str], idx: int, fallback: str) -> str:
    if idx < len(temas):
        return f"{temas[idx].capitalize()} como palanca prioritaria del trimestre."
    return fallback


def armar_resumen_presentacion(slides: List[SlideInfo], tono: str, minutos: int) -> dict[str, List[str] | str]:
    temas = top_temas(slides)
    riesgos = detectar_items(slides, KEYWORDS_RIESGO, max_items=3)
    oportunidades = detectar_items(slides, KEYWORDS_OPORTUNIDAD, max_items=3)

    apertura = {
        "Directo": "Resumen ejecutivo de 2025: foco en resultados, brechas y decisiones inmediatas.",
        "Narrativo": "Cerramos 2025 con aprendizajes claros y una hoja de ruta concreta para 2026.",
        "Data-first": "Partimos de la evidencia: desempeño, riesgos y oportunidades con mayor retorno.",
    }[tono]

    mensajes = [
        mensaje_clave(temas, 0, "Priorizar eficiencia operativa en las iniciativas core."),
        mensaje_clave(temas, 1, "Escalar automatización donde ya hay tracción."),
        mensaje_clave(temas, 2, "Consolidar ejecución transversal entre áreas."),
    ]

    decisiones = [
        "Validar 2 iniciativas con mayor impacto para el próximo trimestre.",
        "Definir sponsor y responsable por iniciativa priorizada.",
        "Acordar revisión ejecutiva con métricas en 30 días.",
    ]

    return {
        "apertura": f"{apertura} (duración sugerida: {minutos} min)",
        "mensajes": mensajes,
        "riesgos": riesgos or ["Sin riesgos explícitos por palabra clave; validar en comité."],
        "oportunidades": oportunidades or ["Sin oportunidades explícitas por palabra clave; complementar con contexto."],
        "decisiones": decisiones,
    }


def app() -> None:
    st.set_page_config(page_title="Executive AI Presenter", page_icon="🧠", layout="wide")
    st.title("🧠 Executive AI Presenter")
    st.caption("Convierte un .pptx en un discurso ejecutivo moderno, con señal clara de uso de IA.")

    base_file = Path("Resumen-Ejecutivo-2025.pptx")

    left, right = st.columns([2, 1])
    with left:
        uploaded = st.file_uploader("Sube una presentación .pptx (opcional)", type=["pptx"])
    with right:
        tono = st.selectbox("Tono", ["Directo", "Narrativo", "Data-first"])
        minutos = st.slider("Duración objetivo (min)", min_value=5, max_value=30, value=12, step=1)
        modo_presentacion = st.toggle("Modo presentación (reunión)", value=True)

    working_file = base_file
    if uploaded is not None:
        working_file = Path(uploaded.name)
        working_file.write_bytes(uploaded.getbuffer())

    if not working_file.exists():
        st.error("No encuentro `Resumen-Ejecutivo-2025.pptx`. Sube el archivo para continuar.")
        st.stop()

    slides = extraer_slides(working_file)
    guion = construir_guion(slides, tono=tono, minutos=minutos)
    resumen = armar_resumen_presentacion(slides, tono=tono, minutos=minutos)

    st.success(f"Archivo analizado: {working_file.name} ({len(slides)} slides)")

    if modo_presentacion:
        st.header("Resumen Ejecutivo para Reunión")
        st.markdown(f"**Apertura sugerida:** {resumen['apertura']}")

        c1, c2 = st.columns(2)
        with c1:
            st.subheader("3 mensajes clave")
            for item in resumen["mensajes"]:
                st.write(f"- {item}")

            st.subheader("Riesgos a vigilar")
            for item in resumen["riesgos"]:
                st.write(f"- {item}")

        with c2:
            st.subheader("Oportunidades de impacto")
            for item in resumen["oportunidades"]:
                st.write(f"- {item}")

            st.subheader("3 decisiones para cerrar")
            for item in resumen["decisiones"]:
                st.write(f"- {item}")

        st.download_button(
            label="Descargar guion ejecutivo (.md)",
            data=guion,
            file_name="guion_ejecutivo_ia.md",
            mime="text/markdown",
        )
    else:
        tab1, tab2, tab3 = st.tabs(["Resumen IA", "Detalle por slide", "Guion para presentar"])

        with tab1:
            temas = top_temas(slides)
            st.subheader("Temas dominantes detectados")
            st.write(" · ".join(temas) if temas else "No se detectaron temas con suficiente texto.")

            st.subheader("Señales de riesgo")
            riesgos = detectar_items(slides, KEYWORDS_RIESGO)
            if riesgos:
                for r in riesgos:
                    st.write(f"- {r}")
            else:
                st.write("No se detectaron riesgos por palabra clave.")

            st.subheader("Señales de oportunidad")
            oportunidades = detectar_items(slides, KEYWORDS_OPORTUNIDAD)
            if oportunidades:
                for o in oportunidades:
                    st.write(f"- {o}")
            else:
                st.write("No se detectaron oportunidades por palabra clave.")

        with tab2:
            for slide in slides:
                with st.expander(f"Slide {slide.number}: {slide.title}"):
                    st.text(slide.body or "(Sin texto)")

        with tab3:
            st.markdown(guion)
            st.download_button(
                label="Descargar guion (.md)",
                data=guion,
                file_name="guion_ejecutivo_ia.md",
                mime="text/markdown",
            )


if __name__ == "__main__":
    app()
