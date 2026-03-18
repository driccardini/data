from pathlib import Path

from pptx import Presentation


OUTPUT_FILE = Path("Presentacion-Ejecutiva-Corporativa-2025-2026.pptx")


def add_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Estrategia Datos & Demanda"
    slide.placeholders[1].text = (
        "Área Demanda & Datos\n"
        "Presentación Corporativa para Dirección\n"
        "Marzo 2026"
    )


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    first = tf.paragraphs[0]
    first.text = bullets[0]
    first.level = 0

    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0


def build() -> Path:
    prs = Presentation()

    add_cover(prs)

    add_bullet_slide(
        prs,
        "Logros 2025",
        [
            "Resultados 2025 y métricas clave",
            "Impacto transversal por áreas de negocio",
            "Avance de transformación tecnológica",
            "Aplicación de Inteligencia Artificial",
            "Prioridades y decisiones para 2026",
        ],
    )

    add_bullet_slide(
        prs,
        "Resultados 2025: Desempeño Consolidado",
        [
            "35 proyectos completados con impacto operativo",
            "8 iniciativas estratégicas en ejecución",
            "+20 tableros manuales automatizados sobre un scope de 120",
            "10 licencias optimizadas con reducción de costos",
            "Consolidación del área como habilitador estratégico",
        ],
    )

    add_bullet_slide(
        prs,
        "Iniciativas por Área e Infraestructura",
        [
            "Administración y Finanzas: automatización de procesos críticos",
            "Integraciones bancarias y tableros de control financiero",
            "Mejoras en trazabilidad y tiempos de respuesta operativa",
            "Estandarización de prácticas de datos para decisiones de gestión",
            "Modernización integral de la base tecnológica",
            "Escalabilidad futura y mayor resiliencia operativa",
            "Optimización de costos de infraestructura y licenciamiento",
            "Base preparada para automatización avanzada e IA",
        ],
    )

    add_bullet_slide(
        prs,
        "Cartera en Progreso: 8 Iniciativas Estratégicas",
        [
            "Foco en cierre de brechas operativas de alto impacto",
            "Automatización de procesos manuales recurrentes",
            "Priorización de quick wins con retorno medible",
            "Gobernanza activa para asegurar ejecución en plazo",
        ],
    )

    add_bullet_slide(
        prs,
        "Inteligencia Artificial Aplicada",
        [
            "IA para priorización de demanda y análisis de oportunidades",
            "Asistencia en generación de reportes ejecutivos",
            "Aceleración de ciclos de análisis y toma de decisiones",
            "Plan 2026: ampliar casos de uso con foco en productividad",
        ],
    )

    add_bullet_slide(
        prs,
        "Riesgos Clave y Mitigación",
        [
            "Riesgo de sobrecarga de cartera: priorización trimestral",
            "Dependencias técnicas: planificación de capacidad",
            "Adopción de nuevas herramientas: gestión del cambio",
            "Seguimiento ejecutivo mensual con KPIs acordados",
        ],
    )

    add_bullet_slide(
        prs,
        "Roadmap 2026",
        [
            "Q2: cierre de iniciativas críticas y estabilización",
            "Q3: escalamiento de automatización e IA en procesos core",
            "Q4: consolidación de beneficios y estandarización regional",
            "Meta anual: mayor eficiencia, control y velocidad de ejecución",
        ],
    )

    add_bullet_slide(
        prs,
        "Decisiones Solicitadas a Dirección",
        [
            "Validar las 3 prioridades estratégicas del próximo trimestre",
            "Asignar sponsors ejecutivos por iniciativa priorizada",
            "Aprobar esquema de seguimiento mensual con tablero único",
            "Confirmar objetivos de impacto para cierre 2026",
        ],
    )

    prs.save(OUTPUT_FILE)
    return OUTPUT_FILE


if __name__ == "__main__":
    output = build()
    print(f"Presentación generada: {output}")