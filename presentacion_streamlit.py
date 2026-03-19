from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List
from html import escape

import streamlit as st
from pptx import Presentation


PPTX_DEFAULT = Path("Presentacion-Ejecutiva-Corporativa-2025-2026.pptx")

ECIPSA_AZUL = "#0B3E53"
ECIPSA_NARANJA = "#F18019"
ECIPSA_CELESTE = "#009AC4"
ECIPSA_GRAFITO = "#2E2F2F"

IMAGENES_ECIPSA = [
    "https://www.ecipsa.com/wp-content/uploads/2021/06/El-Bosque.jpg",
    "https://www.ecipsa.com/wp-content/uploads/2021/06/Valle-Escondido.jpg",
    "https://www.ecipsa.com/wp-content/uploads/2021/06/Valle-Cercano.jpg",
    "https://www.ecipsa.com/wp-content/uploads/2021/06/Tower.jpg",
    "https://www.ecipsa.com/wp-content/uploads/2021/09/foto-milaires.jpg",
]

FONDO_TOWER = "https://www.ecipsa.com/wp-content/uploads/2021/06/Tower.jpg"


@dataclass
class DeckSlide:
    number: int
    title: str
    bullets: List[str]


def extraer_deck(ruta: Path) -> List[DeckSlide]:
    prs = Presentation(str(ruta))
    result: List[DeckSlide] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()

        bullets: List[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or shape == slide.shapes.title:
                continue
            tf = shape.text_frame
            if tf is None:
                continue
            for paragraph in tf.paragraphs:
                text = (paragraph.text or "").strip()
                if text:
                    bullets.append(text)

        if not title:
            title = f"Slide {idx}"

        result.append(DeckSlide(number=idx, title=title, bullets=bullets))

    return result


def estilos() -> None:
    st.markdown(
        f"""
        <style>
            .stApp {{
                background: radial-gradient(circle at 10% 10%, #e9f6fb 0%, #f8fafc 45%, #eef4f8 100%);
                color: {ECIPSA_GRAFITO};
            }}
            .stApp [data-testid="stAppViewContainer"] {{
                position: relative;
                background: transparent;
            }}
            .app-watermark {{
                position: fixed;
                inset: 0;
                background-size: cover;
                background-position: center center;
                opacity: 0.24;
                filter: saturate(0.95) contrast(0.92);
                transform: scale(1.03);
                pointer-events: none;
                z-index: 0;
            }}
            .app-watermark-overlay {{
                position: fixed;
                inset: 0;
                background: linear-gradient(90deg, rgba(248,250,252,0.76) 0%, rgba(248,250,252,0.72) 40%, rgba(248,250,252,0.6) 100%);
                pointer-events: none;
                z-index: 0;
            }}
            .stApp [data-testid="stAppViewContainer"] > .main,
            .stApp header,
            .stApp [data-testid="stToolbar"] {{
                position: relative;
                z-index: 1;
            }}
            .stApp h1,
            .stApp h2,
            .stApp h3,
            .stApp label,
            .stCaption,
            .stMarkdown p,
            .stSlider label {{
                color: {ECIPSA_AZUL} !important;
            }}
            .hero {{
                border: 1px solid {ECIPSA_AZUL};
                border-radius: 20px;
                padding: 20px 24px;
                background: linear-gradient(120deg, {ECIPSA_AZUL} 0%, #134f69 55%, {ECIPSA_CELESTE} 100%);
                color: #f8fafc;
                margin-bottom: 1rem;
                box-shadow: 0 12px 32px rgba(15, 23, 42, 0.22);
            }}
            .hero-title {{
                font-size: 1.7rem;
                font-weight: 700;
                margin: 0.1rem 0 0.5rem 0;
                color: #FFFFFF !important;
            }}
            .hero-sub {{
                font-size: 1rem;
                margin: 0;
                color: #EAF7FC;
            }}
            .deck-wrap {{
                border: 1px solid #cfe4eb;
                border-radius: 20px;
                background: #ffffff;
                padding: 28px 34px;
                min-height: 64vh;
                box-shadow: 0 14px 35px rgba(15, 23, 42, 0.08);
            }}
            .deck-kicker {{
                font-size: 0.9rem;
                color: {ECIPSA_CELESTE};
                letter-spacing: 0.02em;
                margin-bottom: 0.2rem;
                font-weight: 600;
            }}
            .deck-title {{
                font-size: 2.1rem;
                line-height: 1.15;
                font-weight: 700;
                color: {ECIPSA_AZUL};
                margin: 0.3rem 0 1.2rem 0;
            }}
            .deck-bullet {{
                font-size: 1.2rem;
                color: {ECIPSA_GRAFITO};
                margin: 0.45rem 0;
            }}
            .side-card {{
                border: 1px solid #cfe4eb;
                border-radius: 16px;
                background: #ffffff;
                padding: 14px 16px;
                margin-bottom: 0.7rem;
            }}
            .side-panel-slot {{
                min-height: 185px;
                display: flex;
                flex-direction: column;
                justify-content: flex-start;
            }}
            .side-title {{
                font-size: 0.92rem;
                font-weight: 700;
                color: {ECIPSA_AZUL};
                margin: 0;
            }}
            .side-sub {{
                font-size: 0.84rem;
                color: {ECIPSA_GRAFITO};
                margin: 0.2rem 0 0 0;
            }}
            .deck-footer {{
                margin-top: 1.2rem;
                font-size: 0.9rem;
                color: {ECIPSA_GRAFITO};
            }}
            .badge {{
                display: inline-block;
                border: 1px solid #ffd8b5;
                border-radius: 9999px;
                padding: 0.15rem 0.7rem;
                margin-right: 0.45rem;
                font-size: 0.8rem;
                color: #8a4900;
                background: #fff3e8;
            }}
            .stButton > button {{
                border: 1px solid #f4b67f;
                background: #fff7f0;
                color: {ECIPSA_AZUL};
                font-weight: 600;
                min-height: 2.6rem;
                padding: 0.2rem 0.35rem;
                font-size: 1.05rem;
            }}
            .stButton > button:hover {{
                border-color: {ECIPSA_NARANJA};
                color: {ECIPSA_NARANJA};
            }}
            .stProgress > div > div > div > div {{
                background-color: {ECIPSA_NARANJA};
            }}
            .obra-wrap {{
                border: 1px solid #cfe4eb;
                border-radius: 14px;
                background: #ffffff;
                padding: 12px 14px;
                margin: 0.3rem 0 0.8rem 0;
            }}
            .obra-head {{
                color: {ECIPSA_AZUL};
                font-weight: 700;
                font-size: 0.9rem;
                margin-bottom: 0.45rem;
            }}
            .nodes-wrap {{
                display: flex;
                align-items: center;
                gap: 0;
                margin: 0.35rem 0 0.55rem 0;
            }}
            .node {{
                width: 18px;
                height: 18px;
                border-radius: 9999px;
                border: 2px solid #d1d5db;
                background: #f8fafc;
                flex: 0 0 auto;
            }}
            .node.done {{
                border-color: {ECIPSA_NARANJA};
                background: {ECIPSA_NARANJA};
            }}
            .node.current {{
                border-color: {ECIPSA_AZUL};
                background: {ECIPSA_AZUL};
                box-shadow: 0 0 0 3px rgba(11, 62, 83, 0.15);
                animation: pulseNode 1.6s ease-in-out infinite;
            }}
            .node-line {{
                height: 3px;
                background: #dbe3ea;
                flex: 1 1 auto;
                margin: 0 3px;
                border-radius: 99px;
            }}
            .node-line.done {{
                background: {ECIPSA_CELESTE};
            }}
            .tower-svg {{
                margin-top: 6px;
                display: block;
                width: 100%;
                max-width: 320px;
                animation: floatTower 3.8s ease-in-out infinite;
            }}
            .obra-legend {{
                color: {ECIPSA_GRAFITO};
                font-size: 0.82rem;
                margin-top: 0.45rem;
            }}
            @keyframes floatTower {{
                0% {{ transform: translateY(0px); }}
                50% {{ transform: translateY(-2px); }}
                100% {{ transform: translateY(0px); }}
            }}
            @keyframes pulseNode {{
                0% {{ box-shadow: 0 0 0 0 rgba(11, 62, 83, 0.25); }}
                70% {{ box-shadow: 0 0 0 6px rgba(11, 62, 83, 0); }}
                100% {{ box-shadow: 0 0 0 0 rgba(11, 62, 83, 0); }}
            }}
            .metrics-grid {{
                display: grid;
                grid-template-columns: 1fr 1fr;
                gap: 1rem;
                margin: 1.2rem 0 1.4rem 0;
            }}
            .metric-card {{
                border: 1px solid #cfe4eb;
                border-radius: 16px;
                background: #ffffff;
                padding: 1.1rem 1.2rem 1rem 1.2rem;
                box-shadow: 0 4px 14px rgba(15,23,42,0.06);
                display: flex;
                flex-direction: column;
                align-items: center;
                text-align: center;
                gap: 0.25rem;
            }}
            .metric-icon {{
                font-size: 1.6rem;
                line-height: 1;
            }}
            .metric-value {{
                font-size: 2.4rem;
                font-weight: 800;
                color: {ECIPSA_AZUL};
                line-height: 1.1;
            }}
            .metric-label {{
                font-size: 0.9rem;
                color: {ECIPSA_GRAFITO};
                line-height: 1.35;
            }}
            .metric-banner {{
                border-left: 5px solid {ECIPSA_NARANJA};
                border-radius: 0 14px 14px 0;
                background: #fff4ea;
                color: {ECIPSA_GRAFITO};
                font-size: 0.95rem;
                font-weight: 600;
                padding: 0.95rem 1.1rem;
            }}
            .area-grid {{
                display: grid;
                grid-template-columns: 1fr 1fr;
                gap: 0.75rem;
                margin: 0.9rem 0 1rem 0;
            }}
            .area-card {{
                border: 1px solid #cfe4eb;
                border-radius: 14px;
                background: #ffffff;
                padding: 0.85rem 1rem;
                box-shadow: 0 3px 10px rgba(15,23,42,0.05);
            }}
            .area-card-title {{
                font-size: 0.85rem;
                font-weight: 700;
                color: {ECIPSA_AZUL};
                margin-bottom: 0.35rem;
            }}
            .area-card-item {{
                font-size: 0.88rem;
                color: {ECIPSA_GRAFITO};
                line-height: 1.45;
                margin: 0.2rem 0;
            }}
            .area-card-list {{
                margin: 0.25rem 0 0 1.05rem;
                padding: 0;
            }}
            .solution-group {{
                border: 1px solid #b9d9e6;
                border-radius: 18px;
                background: linear-gradient(180deg, #f4fbfe 0%, #eef7fb 100%);
                padding: 1rem 1rem 1.1rem 1rem;
                box-shadow: inset 0 0 0 1px rgba(255,255,255,0.45);
                margin: 0.9rem 0 1rem 0;
            }}
            .solution-group-head {{
                display: flex;
                align-items: center;
                justify-content: space-between;
                gap: 0.8rem;
                margin-bottom: 0.8rem;
                padding-bottom: 0.7rem;
                border-bottom: 1px solid #d8eaf1;
            }}
            .solution-group-title {{
                font-size: 0.96rem;
                font-weight: 800;
                color: {ECIPSA_AZUL};
                margin: 0;
            }}
            .solution-group-sub {{
                font-size: 0.84rem;
                color: {ECIPSA_GRAFITO};
                margin: 0.2rem 0 0 0;
            }}
            .solution-pill {{
                border: 1px solid #b7dbe7;
                border-radius: 999px;
                padding: 0.28rem 0.7rem;
                background: #ffffff;
                color: {ECIPSA_CELESTE};
                font-size: 0.76rem;
                font-weight: 700;
                white-space: nowrap;
            }}
            .solution-grid {{
                display: grid;
                grid-template-columns: repeat(3, minmax(0, 1fr));
                gap: 0.75rem;
            }}
            @media (max-width: 980px) {{
                .solution-grid {{
                    grid-template-columns: 1fr;
                }}
            }}
            .highlight-card {{
                border: 1px solid #cfe4eb;
                border-radius: 16px;
                background: #ffffff;
                padding: 1rem 1.1rem;
                box-shadow: 0 4px 14px rgba(15,23,42,0.06);
            }}
            .highlight-title {{
                font-size: 0.92rem;
                font-weight: 800;
                color: {ECIPSA_AZUL};
                margin: 0 0 0.45rem 0;
            }}
            .highlight-list {{
                margin: 0.2rem 0 0 1.05rem;
                padding: 0;
            }}
            .highlight-item {{
                font-size: 0.9rem;
                color: {ECIPSA_GRAFITO};
                line-height: 1.5;
                margin: 0.22rem 0;
            }}
            .infra-banner {{
                border-left: 4px solid {ECIPSA_CELESTE};
                background: #eaf7fc;
                border-radius: 0 12px 12px 0;
                padding: 0.75rem 1.1rem;
                margin-top: 0.6rem;
            }}
            .infra-banner-title {{
                font-size: 0.88rem;
                font-weight: 700;
                color: {ECIPSA_AZUL};
                margin-bottom: 0.3rem;
            }}
            .infra-banner-item {{
                font-size: 0.88rem;
                color: {ECIPSA_GRAFITO};
                margin: 0.15rem 0;
            }}
        </style>
        """,
        unsafe_allow_html=True,
    )


def imagen_para_slide(numero: int) -> str:
    return IMAGENES_ECIPSA[(numero - 1) % len(IMAGENES_ECIPSA)]


def render_fondo_app(image_url: str) -> None:
    st.markdown(
        f"""
        <div class="app-watermark" style="background-image: url('{image_url}');"></div>
        <div class="app-watermark-overlay"></div>
        """,
        unsafe_allow_html=True,
    )


def icono_para_slide(titulo: str, numero: int) -> str:
    t = titulo.lower()
    if "agenda" in t:
        return "🧭"
    if "resultado" in t or "desempeño" in t:
        return "📈"
    if "impacto" in t:
        return "🏢"
    if "infraestructura" in t:
        return "🖥️"
    if "progreso" in t or "cartera" in t:
        return "🚀"
    if "inteligencia" in t or "ia" in t:
        return "🤖"
    if "riesgo" in t:
        return "⚠️"
    if "roadmap" in t:
        return "🛣️"
    if "decisiones" in t:
        return "✅"
    if numero == 1:
        return "🎯"
    return "📌"


def render_hero(total: int, current: int) -> None:
    st.markdown(
        f"""
        <div class="hero">
            <p class="hero-title">Presentación Ejecutiva Corporativa</p>
            <p class="hero-sub">Formato de exposición moderna para Dirección · Asistida por IA · {current}/{total} slides</p>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_panel_lateral(slide: DeckSlide) -> None:
    if slide.number == 1:
        st.markdown(
            """
            <div class="side-panel-slot">
                <div class="side-card">
                    <p class="side-title">🎯 Propósito de esta presentación</p>
                    <p class="side-sub">Comunicar los logros de 2025, el estado actual de las iniciativas estratégicas y la hoja de ruta para 2026.</p>
                </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        return

    if slide.number == 2:
        st.markdown(
            """
            <div class="side-panel-slot">
                <div class="side-card" style="padding: 1.2rem 1.4rem;">
                    <p class="side-title" style="font-size:1rem;margin-bottom:0.6rem;">🎯 Objetivo</p>
                    <p class="side-sub">Mostrar logros 2025 por área para evidenciar eficientización operativa.</p>
                </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        return

    if slide.number == 4:
        st.markdown(
            """
            <div class="side-panel-slot">
                <div class="side-card" style="padding: 1.2rem 1.4rem;">
                    <p class="side-title" style="font-size:1rem;margin-bottom:0.6rem;">📂 Proyectos</p>
                    <p class="side-sub">Proyectos relacionados a Data & automatizaciones.</p>
                </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        return

    if slide.number == 5:
        st.markdown(
            """
            <div class="side-panel-slot">
                <div class="side-card" style="padding: 1.0rem 1.2rem;margin-bottom:0.8rem;">
                    <p class="side-title" style="font-size:0.88rem;margin-bottom:0.5rem;">📐 Gobierno de Datos</p>
                    <p style="font-size:0.78rem;color:#555;margin:0 0 0.4rem 0;">Estado</p>
                    <div style="background:#e8edf0;border-radius:99px;height:7px;overflow:hidden;">
                        <div style="background:#009AC4;width:20%;height:7px;border-radius:99px;"></div>
                    </div>
                    <p style="font-size:0.82rem;font-weight:700;color:#009AC4;margin:0.25rem 0 0 0;">20% · En inicio</p>
                </div>
                <div class="side-card" style="padding: 1.0rem 1.2rem;">
                    <p class="side-title" style="font-size:0.88rem;margin-bottom:0.5rem;">🤖 Adopción de IA</p>
                    <p style="font-size:0.78rem;color:#555;margin:0 0 0.4rem 0;">Estado</p>
                    <div style="background:#e8edf0;border-radius:99px;height:7px;overflow:hidden;">
                        <div style="background:#F18019;width:10%;height:7px;border-radius:99px;"></div>
                    </div>
                    <p style="font-size:0.82rem;font-weight:700;color:#F18019;margin:0.25rem 0 0 0;">10% · Propuesta en curso</p>
                </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        return

    if slide.number == 6:
        st.markdown(
            """
            <div class="side-panel-slot">
                <div class="side-card" style="padding: 1.1rem 1.2rem;margin-bottom:0.8rem;">
                    <p class="side-title" style="font-size:0.92rem;margin-bottom:0.4rem;">💰 Ahorro estimado</p>
                    <p style="font-size:1.45rem;font-weight:800;color:#F18019;margin:0.1rem 0 0.8rem 0;">~$3.000 USD<span style="font-size:0.85rem;font-weight:500;color:#555;"> / mes</span></p>
                    <p style="font-size:0.78rem;color:#555;margin:0 0 0.3rem 0;">Estado del proyecto</p>
                    <div style="background:#e8edf0;border-radius:99px;height:8px;width:100%;overflow:hidden;">
                        <div style="background:#F18019;width:15%;height:100%;border-radius:99px;"></div>
                    </div>
                    <p style="font-size:0.82rem;font-weight:700;color:#F18019;margin:0.3rem 0 0 0;">15%</p>
                </div>
                <div class="side-card" style="padding: 1.0rem 1.2rem;">
                    <p class="side-title" style="font-size:0.88rem;margin-bottom:0.4rem;">⚙️ Mejora operativa</p>
                    <p class="side-sub" style="margin-bottom:0.5rem;">2 herramientas centralizadas en 1 plataforma unificada.</p>
                    <div style="display:flex;align-items:center;gap:0.5rem;">
                        <span style="font-size:0.82rem;color:#0B3E53;font-weight:700;">🔗 Integración</span>
                        <span style="font-size:0.78rem;background:#e8f4fb;color:#005f7f;border:1px solid #b3ddf2;border-radius:99px;padding:0.1rem 0.55rem;">Dynamics 365</span>
                    </div>
                </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        return

    if slide.number == 101:
        st.markdown(
            """
            <div class="side-panel-slot">
                <div class="side-card" style="padding: 1.1rem 1.2rem;">
                    <p class="side-title" style="font-size:0.9rem;margin-bottom:0.6rem;">&#128101; Composición</p>
                    <div style="display:flex;gap:0.6rem;margin-bottom:0.8rem;">
                        <div style="text-align:center;flex:1;background:#e8f4fb;border-radius:8px;padding:0.5rem 0;">
                            <p style="font-size:1.4rem;font-weight:800;color:#005f7f;margin:0;">2</p>
                            <p style="font-size:0.72rem;color:#555;margin:0;">Internos</p>
                        </div>
                        <div style="text-align:center;flex:1;background:#edf7ed;border-radius:8px;padding:0.5rem 0;">
                            <p style="font-size:1.4rem;font-weight:800;color:#2d6a2d;margin:0;">5</p>
                            <p style="font-size:0.72rem;color:#555;margin:0;">Externos</p>
                        </div>
                    </div>
                    <p style="font-size:0.78rem;color:#555;margin:0 0 0.35rem 0;font-weight:600;">Especialidades</p>
                    <div style="display:flex;flex-wrap:wrap;gap:0.3rem;">
                        <span style="font-size:0.72rem;background:#eef4f8;color:#0B3E53;border:1px solid #c5d9e8;border-radius:99px;padding:0.1rem 0.5rem;">BI &amp; Reporting</span>
                        <span style="font-size:0.72rem;background:#eef4f8;color:#0B3E53;border:1px solid #c5d9e8;border-radius:99px;padding:0.1rem 0.5rem;">Data Engineering</span>
                        <span style="font-size:0.72rem;background:#eef4f8;color:#0B3E53;border:1px solid #c5d9e8;border-radius:99px;padding:0.1rem 0.5rem;">RPA</span>
                        <span style="font-size:0.72rem;background:#eef4f8;color:#0B3E53;border:1px solid #c5d9e8;border-radius:99px;padding:0.1rem 0.5rem;">Gov. Datos</span>
                    </div>
                </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        return

    if slide.number == 100:
        st.markdown(
            """
            <div class="side-panel-slot">
                <div class="side-card" style="padding: 1.2rem 1.4rem;">
                    <p class="side-title" style="font-size:1rem;margin-bottom:0.6rem;">🏗️ Visión</p>
                    <p class="side-sub">Como desarrolladores líderes en real estate, buscamos que la gestión documental, el control de obra y la administración de costos sean procesos integrados, eficientes y a la altura de las mejores prácticas del sector.</p>
                </div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        return

    st.markdown(
        f"""
        <div class="side-panel-slot">
            <div class="side-card">
                <p class="side-title">🧩 Objetivo</p>
                <p class="side-sub">Bajar esta sección a decisiones concretas y medibles para Dirección.</p>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_cover(slide: DeckSlide) -> None:
    st.markdown(
        f"""
        <div class="deck-wrap" style="display:flex;flex-direction:column;justify-content:center;align-items:center;text-align:center;min-height:64vh;">
            <p class="deck-kicker" style="font-size:1rem;letter-spacing:0.08em;">ECIPSA</p>
            <h1 class="deck-title" style="font-size:3rem;margin:0.6rem 0 1.2rem 0;">{escape(slide.title)}</h1>
            <p style="font-size:1.15rem;color:#2E2F2F;max-width:480px;line-height:1.55;">Marzo 2026</p>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_resultados(slide: DeckSlide) -> None:
    areas = [
        (
            "🏦",
            "Administración & Finanzas",
            [
                "6 proyectos implementados",
                "Automatización de rescisiones SAP desde DYN",
                "Tablero de Ingresos Financieros MA",
                "Integración bancaria consolidada (diaria/semanal/mensual)",
                "Notificaciones impositivas ARCA automatizadas",
            ],
        ),
        (
            "🤝",
            "Comercial",
            [
                "5 tableros operacionales",
                "Stock Natania Argentina y Paraguay",
                "Tablero Stock y Precios MA en tiempo real",
                "Tablero Saldo Altoplan",
                "Cupones Natania (Real Time)",
            ],
        ),
        (
            "📞",
            "Ecall",
            [
                "6 reportes críticos",
                "Automatización cupones pendientes de pago",
                "Reportes operativos: Retención, Ventas, NAT, MORA y Auxiliares",
                "Integración directa con Vocalcom",
            ],
        ),
        (
            "🖥️",
            "Modernización Tecnológica",
            [
                "Renovación integral de infraestructura tecnológica",
                "Nueva Infraestructura RPA",
                "Migración Azure con reducción de costos recurrentes",
                "Portal de Datos Ecipsa (+20 reportes y 10 licencias menos)",
            ],
        ),
    ]
    cards_html = "".join(
        f"""<div class='area-card'>
                <div class='area-card-title'>{ico} {titulo}</div>
                <ul class='area-card-list'>
                    {''.join(f"<li class='area-card-item'>{item}</li>" for item in items)}
                </ul>
            </div>"""
        for ico, titulo, items in areas
    )
    st.markdown(
        f"""
        <div class="deck-wrap">
            <h1 class="deck-title" style="font-size:2.3rem;">📈 Iniciativas Implementadas por Área</h1>
            <div class="area-grid">
                {cards_html}
            </div>
            <div class="metric-banner">🏆 Implementaciones 2025 con automatización aplicada y trazabilidad operativa por área</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_logros_por_area(slide: DeckSlide) -> None:
    bullets = [
        "Resultados 2025 y métricas clave",
        "Impacto transversal por áreas de negocio",
        "Avance de transformación tecnológica",
        "Aplicación de Inteligencia Artificial",
        "Prioridades y decisiones para 2026",
    ]

    bullets_html = "".join(f"<li class='area-card-item'>{escape(item)}</li>" for item in bullets)

    st.markdown(
        f"""
        <div class="deck-wrap">
            <h1 class="deck-title" style="font-size:2.3rem;">🏢 Logros por area 2025</h1>
            <div class="area-card">
                <div class="area-card-title">📌 Ejes presentados en la PPTX</div>
                <ul class="area-card-list">
                    {bullets_html}
                </ul>
            </div>
            <div class="metric-banner">🏆 Los avances de 2025 mejoraron eficiencia operativa, visibilidad de gestión y capacidad de respuesta entre áreas</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_areas(slide: DeckSlide) -> None:
    proyectos = [
        (
            "1",
            "Administración (99% · Validación)",
            "Carga renta anticipada en SAP: Automatización de proceso crítico para acelerar reconocimiento de ingresos",
        ),
        (
            "2",
            "Administración (En validación)",
            "Reporte Aportes y Saldos MA: Sincronización automática con DYN para visibilidad financiera en tiempo real",
        ),
        (
            "3",
            "Créditos & Cobranzas (30%)",
            "Cobranzas manuales: Automatización de gestión para reducir tiempos de recuperación y liberar recursos",
        ),
        (
            "4",
            "Data (30%)",
            "Reportes manuales: Eliminación de procesos repetitivos mediante automatización inteligente",
        ),
        (
            "5",
            "Tax (En desarrollo)",
            "Control de impuestos: Automatización de cálculos y controles para reducir riesgo fiscal",
        ),
        (
            "6",
            "IT - IA (Implementación)",
            "Lectura de comprobantes con IA: Bot WhatsApp que reconoce transferencias bancarias para agilizar cobranzas manuales",
        ),
        (
            "7",
            "Comercial (En curso)",
            "Cartera de clientes: Vista 360° centralizada del cliente para mejorar experiencia y estrategia comercial",
        ),
    ]
    cards_html = "".join(
        f"""<div class='area-card'>
                <div class='area-card-title'>#{num} {titulo}</div>
                <ul class='area-card-list'>
                    <li class='area-card-item'>{descripcion}</li>
                </ul>
            </div>"""
        for num, titulo, descripcion in proyectos
    )
    st.markdown(
        f"""
        <div class="deck-wrap">
            <h1 class="deck-title" style="font-size:2.3rem;">🚧 Proyectos en Desarrollo Activo</h1>
            <div class="area-grid">
                {cards_html}
            </div>
            <div class="metric-banner">🏆 Ocho iniciativas estratégicas continúan avanzando para completar la transformación digital del 2025, con foco en automatización de procesos manuales e inteligencia artificial aplicada.</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_genesys(slide: DeckSlide) -> None:
    st.markdown(
        """
        <div class="deck-wrap">
            <h1 class="deck-title" style="font-size:2.3rem;">📡 Genesys CX</h1>
            <div class="solution-group">
                <div class="solution-group-head">
                    <div>
                        <p class="solution-group-title">Reemplazo de Vocalcom & WISE · Canales de voz y digital</p>
                        <p class="solution-group-sub">Unifica WhatsApp, Instagram, Facebook y voz en una sola plataforma, con IA incorporada y ahorro mensual operativo.</p>
                    </div>
                    <span class="solution-pill">Discovery activo</span>
                </div>
                <div class="solution-grid">
                    <div class="area-card">
                        <div class="area-card-title">🎙️ Canal de voz</div>
                        <ul class="area-card-list">
                            <li class="area-card-item">Revisión y mejora de flujos IVR.</li>
                            <li class="area-card-item">Incorporación de IA en la atención por voz.</li>
                        </ul>
                    </div>
                    <div class="area-card">
                        <div class="area-card-title">💬 WhatsApp</div>
                        <ul class="area-card-list">
                            <li class="area-card-item">Relevamiento de bots actuales.</li>
                            <li class="area-card-item">Revisión y mejora de procesos de atención digital.</li>
                        </ul>
                    </div>
                    <div class="area-card">
                        <div class="area-card-title">🎓 Capacitación BVS</div>
                        <ul class="area-card-list">
                            <li class="area-card-item">Reuniones en vivo con la consultora implementadora.</li>
                            <li class="area-card-item">Demos y documentación de los canales.</li>
                        </ul>
                    </div>
                </div>
            </div>
            <div class="metric-banner">🏆 Genesys CX reemplaza Vocalcom y WISE generando ahorro mensual y habilitando IA en todos los canales de contacto.</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_herramienta_obra(slide: DeckSlide) -> None:
    st.markdown(
        """
        <div class="deck-wrap">
            <h1 class="deck-title" style="font-size:2.3rem;">🏗️ ECIPSA Obra Hub</h1>
            <div class="solution-group">
                <div class="solution-group-head">
                    <div>
                        <p class="solution-group-title">Portal único de gestión de obra</p>
                        <p class="solution-group-sub">Tres capacidades incluidas dentro de una misma solución para proveedores, seguimiento y control interno.</p>
                    </div>
                    <span class="solution-pill">3 módulos integrados</span>
                </div>
                <div class="solution-grid" style="align-items:end;">
                    <div class="area-card">
                        <div class="area-card-title">📄 1. Gestión documental</div>
                        <ul class="area-card-list">
                            <li class="area-card-item">Herramienta para que los proveedores puedan interactuar con Ecipsa subiendo planos y documentación necesaria para trabajar.</li>
                        </ul>
                        <div style='margin-top:auto;'>
                            <span style='font-size:0.85rem;font-weight:700;color:#009AC4;'>70% · Relevado, con propuestas comerciales de 3 proveedores</span>
                            <div style='background:#e8edf0;border-radius:99px;height:7px;overflow:hidden;margin-top:0.18rem;'>
                                <div style='background:#009AC4;width:70%;height:7px;border-radius:99px;'></div>
                            </div>
                        </div>
                    </div>
                    <div class="area-card">
                        <div class="area-card-title">🏗️ 2. Dirección de obra</div>
                        <ul class="area-card-list">
                            <li class="area-card-item">Herramienta para poder darle seguimiento a la obra constructiva y que los proveedores puedan ir actualizando el avance de la misma.</li>
                        </ul>
                        <div style='margin-top:auto;'>
                            <span style='font-size:0.85rem;font-weight:700;color:#F18019;'>15% · En análisis junto con el equipo de obra de Córdoba</span>
                            <div style='background:#e8edf0;border-radius:99px;height:7px;overflow:hidden;margin-top:0.18rem;'>
                                <div style='background:#F18019;width:15%;height:7px;border-radius:99px;'></div>
                            </div>
                        </div>
                    </div>
                    <div class="area-card">
                        <div class="area-card-title">💰 3. Gestión de costos de obra</div>
                        <ul class="area-card-list">
                            <li class="area-card-item">Herramienta interna para poder calcular los costos correctos de la obra.</li>
                        </ul>
                        <div style='margin-top:auto;'>
                            <span style='font-size:0.85rem;font-weight:700;color:#0B3E53;'>20% · En análisis</span>
                            <div style='background:#e8edf0;border-radius:99px;height:7px;overflow:hidden;margin-top:0.18rem;'>
                                <div style='background:#0B3E53;width:20%;height:7px;border-radius:99px;'></div>
                            </div>
                        </div>
                    </div>
                </div>
            </div>
            <div class="metric-banner">🏆 Estas tres herramientas van a converger en una sola solución: ECIPSA Obra Hub.</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_gobierno_ia(slide: DeckSlide) -> None:
    st.markdown(
        """
        <div class="deck-wrap">
            <h1 class="deck-title" style="font-size:2.1rem;">🏛️ Gobierno de Datos &amp; Adopción de IA</h1>
            <div style="display:grid;grid-template-columns:1fr 1fr;gap:1.2rem;margin-top:0.6rem;">
                <div class="highlight-card" style="display:flex;flex-direction:column;gap:0.5rem;">
                    <div class="highlight-title" style="font-size:1rem;">📐 Gobierno de Datos</div>
                    <p style="font-size:0.88rem;color:#444;margin:0 0 0.5rem 0;">Sentando las bases de trazabilidad extremo a extremo sobre los datos de ECIPSA.</p>
                    <div style="display:flex;flex-wrap:wrap;gap:0.4rem;margin-bottom:0.5rem;">
                        <span style="font-size:0.78rem;background:#e8f4fb;color:#005f7f;border:1px solid #b3ddf2;border-radius:99px;padding:0.15rem 0.6rem;font-weight:600;">PureView</span>
                        <span style="font-size:0.78rem;background:#eef4f8;color:#0B3E53;border:1px solid #c5d9e8;border-radius:99px;padding:0.15rem 0.6rem;font-weight:600;">DWH</span>
                        <span style="font-size:0.78rem;background:#fff3e8;color:#8a4900;border:1px solid #ffd8b5;border-radius:99px;padding:0.15rem 0.6rem;font-weight:600;">Dynamics 365</span>
                        <span style="font-size:0.78rem;background:#f0f7ff;color:#1a5fa8;border:1px solid #c2d9f5;border-radius:99px;padding:0.15rem 0.6rem;font-weight:600;">Power BI</span>
                    </div>
                    <ul class="highlight-list">
                        <li class="highlight-item">Integración de PureView con el Data Warehouse, Dynamics y Power BI.</li>
                        <li class="highlight-item">Trazabilidad completa desde el origen del dato hasta su consumo en reportes.</li>
                        <li class="highlight-item">Base para certificación de calidad de datos en procesos críticos.</li>
                    </ul>
                </div>
                <div class="highlight-card" style="display:flex;flex-direction:column;gap:0.5rem;">
                    <div class="highlight-title" style="font-size:1rem;">🤖 Adopción de IA</div>
                    <p style="font-size:0.88rem;color:#444;margin:0 0 0.5rem 0;">Propuesta para trasladar el uso de IA al día a día de los colaboradores en áreas clave.</p>
                    <div style="display:flex;flex-wrap:wrap;gap:0.4rem;margin-bottom:0.5rem;">
                        <span style="font-size:0.78rem;background:#eef4f8;color:#0B3E53;border:1px solid #c5d9e8;border-radius:99px;padding:0.15rem 0.6rem;">📒 Contabilidad</span>
                        <span style="font-size:0.78rem;background:#eef4f8;color:#0B3E53;border:1px solid #c5d9e8;border-radius:99px;padding:0.15rem 0.6rem;">💰 Finanzas</span>
                        <span style="font-size:0.78rem;background:#eef4f8;color:#0B3E53;border:1px solid #c5d9e8;border-radius:99px;padding:0.15rem 0.6rem;">📊 Control de Gestión</span>
                        <span style="font-size:0.78rem;background:#eef4f8;color:#0B3E53;border:1px solid #c5d9e8;border-radius:99px;padding:0.15rem 0.6rem;">👥 RRHH</span>
                        <span style="font-size:0.78rem;background:#eef4f8;color:#0B3E53;border:1px solid #c5d9e8;border-radius:99px;padding:0.15rem 0.6rem;">🔔 Cobranzas</span>
                    </div>
                    <ul class="highlight-list">
                        <li class="highlight-item">Definir propuesta de adopción con foco en productividad operativa.</li>
                        <li class="highlight-item">Capacitación y acompañamiento por área, comenzando por las candidatas identificadas.</li>
                        <li class="highlight-item">Métricas de adopción y seguimiento de impacto por equipo.</li>
                    </ul>
                </div>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_equipo(slide: DeckSlide) -> None:
    _seniority_order = {"Senior": 0, "Semi Senior": 1, "Junior": 2}

    _internos = sorted([
        ("Romina Meade",       "Responsable de Reporting",    "Senior",      ["Portal de Datos", "Genesys", "Portal Proveedores", "Autom. Ecall", "Autom. Bancarias", "Integr. SAP", "Costo de Obra"]),
        ("Agust\u00edn Contreras", "Analista Senior de Datos",    "Senior",      ["Portal de Datos", "Integraciones SAP", "Gob. de Datos"]),
    ], key=lambda x: _seniority_order.get(x[2], 9))

    _externos = sorted([
        ("Florencia Olmedo",   "Desarrolladora BI",            "Semi Senior", ["Tableros MA", "Tableros Comercial", "Portal Licitaciones", "Tableros Ecall", "Tableros MKT"]),
        ("Mateo Cuenca",       "Ingeniero de Datos",           "Semi Senior", ["Integr. SAP", "Integr. ARCA", "Integr. NITRO", "WhatsApp Business"]),
        ("Mat\u00edas Cattaneo",   "Desarrollador BI & Gob. Datos", "Junior",      ["Tableros Ecall", "Gob. Datos \u00b7 Pureview", "Cat\u00e1logos Power BI"]),
        ("Maximiliano Rivera", "Desarrollador BI",             "Junior",      ["Tableros Power BI", "Reporte Postventa"]),
        ("Alan Riquelmes",     "Automatizaciones RPA",         "Junior",      ["Extractos Bancarios", "Chatbot Cobranzas", "Chatbot C-Level"]),
    ], key=lambda x: _seniority_order.get(x[2], 9))

    def _card(nombre, rol, seniority, proyectos, _bg):
        pills = "".join(
            f"<span style='font-size:0.66rem;background:#eef4f8;color:#0B3E53;border-radius:4px;padding:0.07rem 0.36rem;'>{p}</span>"
            for p in proyectos
        )
        sen_col = {"Senior": "#c07000", "Semi Senior": "#005f7f", "Junior": "#3a7a3a"}.get(seniority, "#555")
        sen_bg  = {"Senior": "#fff3e0", "Semi Senior": "#e8f4fb", "Junior": "#edf7ed"}.get(seniority, "#eee")
        return f"""
        <div class="highlight-card" style="padding:0.75rem 0.9rem;">
            <span style="font-size:0.88rem;font-weight:800;color:#0B3E53;display:block;margin-bottom:0.18rem;">{nombre}</span>
            <p style="font-size:0.73rem;color:#666;margin:0 0 0.35rem 0;">{rol}</p>
            <span style="font-size:0.67rem;background:{sen_bg};color:{sen_col};border-radius:99px;padding:0.07rem 0.42rem;font-weight:600;">{seniority}</span>
            <div style="display:flex;flex-wrap:wrap;gap:0.25rem;margin-top:0.4rem;">{pills}</div>
        </div>"""

    internos_html = "".join(_card(*p, "") for p in _internos)
    externos_html = "".join(_card(*p, "") for p in _externos)

    st.markdown(
        f"""
        <div class="deck-wrap">
            <h1 class="deck-title" style="font-size:2rem;">\U0001f465 Equipo Data &amp; Automatizaciones</h1>
            <div style="display:flex;flex-direction:column;gap:1.2rem;margin-top:0.5rem;">
                <div>
                    <p style="font-size:0.78rem;font-weight:800;color:#005f7f;text-transform:uppercase;letter-spacing:0.07em;margin:0 0 0.55rem 0;border-bottom:2px solid #b3ddf2;padding-bottom:0.25rem;">\U0001f7e6 Planta Interna &mdash; 2 colaboradores</p>
                    <div style="display:flex;flex-direction:column;gap:0.6rem;">{internos_html}</div>
                </div>
                <div>
                    <p style="font-size:0.78rem;font-weight:800;color:#2d6a2d;text-transform:uppercase;letter-spacing:0.07em;margin:1.2rem 0 0.55rem 0;border-bottom:2px solid #b5d8b5;padding-bottom:0.25rem;">\U0001f7e9 Equipo Externo &mdash; 5 colaboradores</p>
                    <div style="display:grid;grid-template-columns:1fr 1fr;gap:0.6rem;">{externos_html}</div>
                </div>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_portal_proveedores(slide: DeckSlide) -> None:
    st.markdown(
        """
        <div class="deck-wrap">
            <h1 class="deck-title" style="font-size:2.3rem;">🤝 Nuevo portal de proveedores</h1>
            <div class="area-grid">
                <div class="highlight-card">
                    <div class="highlight-title">📍 Status actual</div>
                    <ul class="highlight-list">
                        <li class="highlight-item">Pruebas piloto con proveedores en curso.</li>
                        <li class="highlight-item">Estrategia de go live en definición junto al equipo de Comunicaciones.</li>
                        <li class="highlight-item">Validación interna de los procesos de aprobación.</li>
                    </ul>
                </div>
                <div class="highlight-card">
                    <div class="highlight-title">⏳ Pendientes clave</div>
                    <ul class="highlight-list">
                        <li class="highlight-item">Terminar validaciones de comprobantes de Paraguay.</li>
                        <li class="highlight-item">Completar el proceso automático de alta de proveedores.</li>
                    </ul>
                </div>
                <div class="highlight-card" style="grid-column: 1 / -1;">
                    <div class="highlight-title">🏆 Logros esperados</div>
                    <ul class="highlight-list">
                        <li class="highlight-item">Prescindir de estudios contables que hoy cargan facturas en Argentina y próximamente en Paraguay.</li>
                        <li class="highlight-item">Integración completa con SAP.</li>
                        <li class="highlight-item">Mejor gestión integral de proveedores.</li>
                    </ul>
                </div>
            </div>
            <div class="metric-banner">🏆 Implementación orientada a digitalizar punta a punta la relación con proveedores, con mayor control, trazabilidad e integración operativa.</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_roadmap_2026(slide: DeckSlide) -> None:
    st.markdown(
        f"""
        <div class="deck-wrap" style="position:relative;overflow:hidden;min-height:60vh;">
            <img src='https://www.ecipsa.com/wp-content/uploads/2021/06/Tower.jpg' alt='Fondo Roadmap' style='position:absolute;top:0;left:0;width:100%;height:100%;object-fit:cover;opacity:0.13;z-index:0;pointer-events:none;'>
            <div style="position:relative;z-index:1;">
                <h1 class="deck-title" style="font-size:2.3rem;">🛣️ Roadmap 2026</h1>
                <ul style='font-size:1.13rem;color:#2E2F2F;max-width:700px;margin:2.2rem auto 1.2rem auto;line-height:1.55;'>
                    <li>🌐 <b>Portal de Proveedores:</b> Finalizar la digitalización y lanzar la v2, incluyendo alta de proveedores y expansión a Paraguay.</li>
                    <li>🧑‍💻 <b>Portal de Clientes:</b> Transformar el portal de licitaciones en el nuevo portal de clientes y eliminar la app MiNatania.</li>
                    <li>🏗️ <b>Gestión de Obra:</b> Integrar tecnología con Dyn y SAP para mayor control y eficiencia en obra.</li>
                    <li>📞 <b>Genesys:</b> Implementar Genesys en todos los canales de contacto para mejorar la experiencia del cliente.</li>
                    <li>🤖 <b>Automatización (RPA):</b> Analizar e implementar herramientas de RPA para ahorrar tiempo y mejorar procesos.</li>
                    <li>📊 <b>Administración:</b> Relevar y optimizar procesos administrativos para reducir tareas manuales.</li>
                    <li>👥 <b>Recursos Humanos:</b> Colaborar para automatizar y mejorar procesos clave de HR.</li>
                    <li>💻 <b>Asset Management:</b> Automatizar la baja de equipos junto con la baja del colaborador, integrado con Bizneo, IT y Arquitectura.</li>
                    <li>💳 <b>Control de Gestión:</b> Mejorar procesos de pago y rendición de software contratado con tarjetas corporativas.</li>
                    <li>🧾 <b>SAP:</b> Automatizar controles para detectar desvíos de presupuesto y ejecución.</li>
                    <li>📄 <b>Contratos Inteligentes:</b> Evaluar herramientas con IA para gestión de contratos junto a Compras y Legales.</li>
                    <li>📈 <b>Planificación Financiera:</b> Buscar e implementar soluciones innovadoras junto a Finanzas y Nuevos Negocios.</li>
                </ul>
                <div class="metric-banner" style="margin-top:2.2rem;">🚀 <b>2026:</b> Consolidación digital, eficiencia operativa y experiencia superior para todos los actores de ECIPSA.</div>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

def render_slide(slide: DeckSlide, total: int) -> None:
    if slide.number == 1:
        render_cover(slide)
        return

    if slide.number == 2:
        render_logros_por_area(slide)
        return

    if slide.number == 3:
        render_resultados(slide)
        return

    if slide.number == 4:
        render_areas(slide)
        return

    if slide.number == 5:
        render_gobierno_ia(slide)
        return

    if slide.number == 6:
        render_genesys(slide)
        return

    if slide.number == 99:
        render_roadmap_2026(slide)
        return

    if slide.number == 101:
        render_equipo(slide)
        return

    if slide.number == 100:
        render_herramienta_obra(slide)
        return

    if slide.number == 99:
        render_roadmap_2026(slide)
        return

    tam_bullet = "1.35rem"
    tam_titulo = "2.5rem"
    icono = icono_para_slide(slide.title, slide.number)

    bullets_html = "".join(
        f"<li class='deck-bullet' style='font-size:{tam_bullet};'>{escape(b)}</li>" for b in slide.bullets[:8]
    )

    st.markdown(
        f"""
        <div class="deck-wrap">
            <h1 class="deck-title" style="font-size:{tam_titulo};">{icono} {escape(slide.title)}</h1>
            <ul style="padding-left: 1.4rem; margin-top: 0.2rem;">
                {bullets_html}
            </ul>
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_navegacion_obra(actual: int, total: int) -> None:
    """Edificio que crece piso a piso: cada slide = un piso nuevo."""
    FLOOR_H = max(8, min(18, 140 // total))
    BLD_X = 50
    BLD_W = 80
    GROUND_Y = 168
    NUM_WINS = 4
    WIN_W = 9
    win_h = max(4, FLOOR_H - 6)
    spacing = (BLD_W - NUM_WINS * WIN_W) / (NUM_WINS + 1)

    parts: list = []

    # ── floors ────────────────────────────────────────────────────────────
    for i in range(actual):
        floor_y = GROUND_Y - (i + 1) * FLOOR_H
        is_current = (i == actual - 1)
        fill = "#0f5272" if is_current else "#0B3E53"

        # floor slab
        if is_current and actual > 1:
            parts.append(
                f"<rect x='{BLD_X}' y='{floor_y}' width='{BLD_W}' height='{FLOOR_H - 1}' fill='{fill}' rx='2'>"
                "<animate attributeName='opacity' values='0;1' dur='0.5s' fill='freeze'/></rect>"
            )
        else:
            parts.append(
                f"<rect x='{BLD_X}' y='{floor_y}' width='{BLD_W}' height='{FLOOR_H - 1}' fill='{fill}' rx='2'/>"
            )

        # windows
        win_y = floor_y + 2
        for j in range(NUM_WINS):
            wx = round(BLD_X + spacing + j * (WIN_W + spacing))
            if is_current:
                parts.append(
                    f"<rect x='{wx}' y='{win_y}' width='{WIN_W}' height='{win_h}' rx='1.5' fill='#FFFDE7' stroke='#F18019' stroke-width='0.8'>"
                    f"<animate attributeName='opacity' values='0.5;1;0.7;1' dur='1.8s' begin='{j * 0.15:.2f}s' repeatCount='indefinite'/></rect>"
                )
            else:
                parts.append(
                    f"<rect x='{wx}' y='{win_y}' width='{WIN_W}' height='{win_h}' rx='1.5' fill='#BFD9E8' stroke='#9BC9DF' stroke-width='0.6'/>"
                )

    # ── crane when building is in progress ────────────────────────────────
    if actual < total:
        crane_top = GROUND_Y - actual * FLOOR_H
        arm_h = max(16, int(FLOOR_H * 1.5))
        mast_top = crane_top - arm_h
        arm_len = 26
        wire_end_y = mast_top + arm_h // 2
        ball_cx = BLD_X + BLD_W - 2  # right side of building
        parts += [
            f"<rect x='88' y='{mast_top}' width='4' height='{arm_h}' fill='{ECIPSA_NARANJA}' rx='1'/>",
            f"<rect x='88' y='{mast_top}' width='{arm_len}' height='3' fill='{ECIPSA_NARANJA}' rx='1'/>",
            f"<line x1='{88 + arm_len - 1}' y1='{mast_top + 3}' x2='{88 + arm_len - 1}' y2='{wire_end_y}' stroke='{ECIPSA_NARANJA}' stroke-width='1' stroke-dasharray='2,2'/>",
            f"<circle cx='{88 + arm_len - 1}' cy='{wire_end_y + 4}' r='3.5' fill='{ECIPSA_NARANJA}'>"
            "<animate attributeName='opacity' values='0.3;1;0.3' dur='0.9s' repeatCount='indefinite'/>"
            "</circle>",
        ]

    # ── flag when complete ────────────────────────────────────────────────
    if actual == total:
        top_y = GROUND_Y - total * FLOOR_H
        parts += [
            f"<rect x='88' y='{top_y - 18}' width='2.5' height='18' fill='{ECIPSA_GRAFITO}'/>",
            f"<polygon points='90.5,{top_y - 18} 101,{top_y - 13} 90.5,{top_y - 8}' fill='{ECIPSA_NARANJA}'/>",
        ]

    # ── ground line ───────────────────────────────────────────────────────
    parts.append(f"<rect x='0' y='{GROUND_Y}' width='180' height='5' fill='#cfe4eb' rx='2'/>")

    svg = (
        "<svg viewBox='0 0 180 200' xmlns='http://www.w3.org/2000/svg' "
        "style='width:100%;max-width:200px;display:block;margin:0 auto;'>"
        + "".join(parts)
        + "</svg>"
    )

    icono = "🏗️" if actual < total else "🏢"
    piso_txt = "Planta baja" if actual == 1 else f"Piso {actual - 1}"

    st.markdown(
        f"""
        <div class="obra-wrap">
            <div class="obra-head">{icono} {piso_txt} &mdash; slide {actual} de {total}</div>
            {svg}
        </div>
        """,
        unsafe_allow_html=True,
    )


def app() -> None:
    st.set_page_config(page_title="Presentación Ejecutiva", page_icon="📊", layout="wide")
    estilos()

    pptx_path = PPTX_DEFAULT

    if not pptx_path.exists():
        st.error("No se encontró la presentación corporativa base en la carpeta del proyecto.")
        st.stop()

    _all = [s for s in extraer_deck(pptx_path) if s.number != 2]
    # Insertar Roadmap 2026 como piso 7 (después de ECIPSA Obra Hub)
    piso_roadmap = DeckSlide(
        number=99,
        title="Roadmap 2026",
        bullets=[],
    )
    piso_herramienta_obra = DeckSlide(
        number=100,
        title="ECIPSA Obra Hub",
        bullets=["Mil Aires (En análisis)", "Herramienta dirección de obra: Evaluación de soluciones SAAS vs desarrollo a medida"],
    )
    piso_equipo = DeckSlide(number=101, title="Equipo", bullets=[])
    _by_num = {s.number: s for s in _all}
    _front_order = [1, 3, 4, 5, 6]
    _base = [_by_num[n] for n in _front_order if n in _by_num]
    slides = [_base[0], piso_equipo] + _base[1:]
    slides.append(piso_herramienta_obra)
    slides.append(piso_roadmap)
    slides += [s for s in _all if s.number not in _front_order]
    total = len(slides)

    if "slide_idx" not in st.session_state:
        st.session_state.slide_idx = 0

    # Procesar navegacion ANTES de renderizar para que la slide correcta aparezca al primer click
    if st.session_state.get("_nav") == "prev":
        st.session_state.slide_idx = max(0, st.session_state.slide_idx - 1)
        del st.session_state["_nav"]
    elif st.session_state.get("_nav") == "next":
        st.session_state.slide_idx = min(total - 1, st.session_state.slide_idx + 1)
        del st.session_state["_nav"]

    render_fondo_app(FONDO_TOWER)


    content_l, content_r = st.columns([3.2, 1.2])
    with content_l:
        render_slide(slides[st.session_state.slide_idx], total)
    with content_r:
        render_panel_lateral(slides[st.session_state.slide_idx])
        render_navegacion_obra(st.session_state.slide_idx + 1, total)
        nav_l, _, nav_r = st.columns([1, 1.2, 1])
        with nav_l:
            if st.button("◀", use_container_width=True):
                st.session_state["_nav"] = "prev"
                st.rerun()
        with nav_r:
            if st.button("▶", use_container_width=True):
                st.session_state["_nav"] = "next"
                st.rerun()

    st.progress((st.session_state.slide_idx + 1) / total)


if __name__ == "__main__":
    app()
